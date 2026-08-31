# -*- coding: utf-8 -*-
"""HazardBot 덱 — design.md 를 python-pptx 도형으로 옮긴 프리미티브·컴포넌트.

1280x720 px 아트보드를 그대로 좌표로 쓴다. 1px = 9525 EMU = 0.75pt.
여기에는 '어떻게 보이는가'만 있고 '무엇을 말하는가'는 build_deck.py 에 있다.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
import unicodedata

PX = 9525


def E(v):
    return Emu(int(round(v * PX)))


def P(px):
    return Pt(px * 0.75)


# ── §5 색 토큰 ────────────────────────────────────────────────────────────────
GROUND    = RGBColor(0xF5, 0xF5, 0xF7)   # canvas-parchment
CANVAS    = RGBColor(0xFF, 0xFF, 0xFF)   # canvas
PEARL     = RGBColor(0xFA, 0xFA, 0xFC)   # surface-pearl
TILE1     = RGBColor(0x27, 0x27, 0x29)   # surface-tile-1
TILE2     = RGBColor(0x2A, 0x2A, 0x2C)   # surface-tile-2
TILE3     = RGBColor(0x25, 0x25, 0x27)   # surface-tile-3
BLACK     = RGBColor(0x00, 0x00, 0x00)   # surface-black
INK       = RGBColor(0x1D, 0x1D, 0x1F)   # ink
INK_DIM   = RGBColor(0x33, 0x33, 0x33)   # ink-muted-80
INK_FNT   = RGBColor(0x7A, 0x7A, 0x7A)   # ink-muted-48
ON_DARK   = RGBColor(0xFF, 0xFF, 0xFF)   # body-on-dark
ON_DARK_M = RGBColor(0xCC, 0xCC, 0xCC)   # body-muted
RULE      = RGBColor(0xE0, 0xE0, 0xE0)   # hairline
RULE_SFT  = RGBColor(0xF0, 0xF0, 0xF0)   # divider-soft
RULE_DARK = RGBColor(0x3A, 0x3A, 0x3C)   # hairline-on-dark
ACCENT    = RGBColor(0x00, 0x66, 0xCC)   # primary
ACCENT_F  = RGBColor(0x00, 0x71, 0xE3)   # primary-focus
ACCENT_D  = RGBColor(0x29, 0x97, 0xFF)   # primary-on-dark
WASH      = RGBColor(0xEB, 0xF3, 0xFB)   # primary-wash
HAZARD    = RGBColor(0xD7, 0x00, 0x15)   # status-hazard
OXIDIZER  = RGBColor(0xB2, 0x50, 0x00)   # status-oxidizer
CLEAR     = RGBColor(0x24, 0x8A, 0x3D)   # status-clear

# 상태 14% 틴트를 흰 면(#ffffff) / 파치먼트(#f5f5f7) 위에 미리 합성.
# PowerPoint 도형 채움에는 알파가 없어서 계산해 넣는다.
def _tint(c, base, a=0.14):
    return RGBColor(*[int(round(c[i] * a + base[i] * (1 - a))) for i in range(3)])


_W = (0xFF, 0xFF, 0xFF)
_G = (0xF5, 0xF5, 0xF7)
HAZARD_BG   = _tint((0xD7, 0x00, 0x15), _W)
OXIDIZER_BG = _tint((0xB2, 0x50, 0x00), _W)
CLEAR_BG    = _tint((0x24, 0x8A, 0x3D), _W)
HAZARD_BG_G   = _tint((0xD7, 0x00, 0x15), _G)
OXIDIZER_BG_G = _tint((0xB2, 0x50, 0x00), _G)
CLEAR_BG_G    = _tint((0x24, 0x8A, 0x3D), _G)

STATUS = {"hazard":   (HAZARD,   HAZARD_BG,   HAZARD_BG_G),
          "oxidizer": (OXIDIZER, OXIDIZER_BG, OXIDIZER_BG_G),
          "clear":    (CLEAR,    CLEAR_BG,    CLEAR_BG_G)}

# ── §3 프레임 ─────────────────────────────────────────────────────────────────
BODY_X, BODY_Y, BODY_W, BODY_H = 64, 176, 1152, 488
BODY_B = BODY_Y + BODY_H          # 664
GAP = 24                          # §7 거터 = 행간격 = 24


def span(n):
    """§7 스팬 폭: 98n - 24"""
    return 98 * n - 24


# ── §4 Pretendard: 600은 별도 패밀리, 700은 Bold 비트 ────────────────────────
def face(weight):
    return ("Pretendard SemiBold", False) if weight == 600 else \
           ("Pretendard", weight >= 700)


# ── 문서·슬라이드 상태 ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = E(1280)
prs.slide_height = E(720)
SH = None
WARNINGS = []


def new_slide():
    global SH
    s = prs.slides.add_slide(prs.slide_layouts[6])
    SH = s.shapes
    return s


def save(path):
    prs.save(path)


# ── 텍스트 폭 추정 — 상자 밖으로 넘치는 글을 빌드 시점에 잡아낸다 ────────────
def _cw(ch):
    o = ord(ch)
    if 0xAC00 <= o <= 0xD7A3 or 0x3130 <= o <= 0x318F or 0x4E00 <= o <= 0x9FFF:
        return 1.00
    if ch in "—…":
        return 1.00
    if ch == " ":
        return 0.26
    if ch in "·":
        return 0.40
    if ch.isdigit():
        return 0.56
    if ch.isupper():
        return 0.68
    if ch.islower():
        return 0.55
    if unicodedata.east_asian_width(ch) in ("W", "F"):
        return 1.00
    return 0.34


def text_w(s, size, weight=400, spc=0.0):
    f = 1.04 if weight >= 600 else 1.0
    return sum(_cw(c) for c in s) * size * f + len(s) * size * spc


def est_lines(text, size, width, weight=400, spc=0.0):
    """어절 단위(keep-all)로 줄바꿈했을 때의 줄 수를 보수적으로 추정."""
    lines, cur = 1, 0.0
    sp = text_w(" ", size, weight, spc)
    for word in text.split(" "):
        w = text_w(word, size, weight, spc)
        if cur == 0:
            cur = w
        elif cur + sp + w <= width:
            cur += sp + w
        else:
            lines += 1
            cur = w
    return lines


def check(tag, text, size, width, height, lh, weight=400, spc=0.0):
    n = est_lines(text, size, width, weight, spc)
    need = n * size * lh
    if need > height + 0.5:
        WARNINGS.append(
            f"[overflow] {tag}: {n}줄 x {size}px x {lh} = {need:.0f}px > {height}px  "
            f"| {text[:44]}")
    return n


# ── 프리미티브 ────────────────────────────────────────────────────────────────
def no_shadow(shape):
    spPr = shape._element.spPr
    if spPr.find(qn('a:effectLst')) is None:
        spPr.append(spPr.makeelement(qn('a:effectLst'), {}))


def rect(x, y, w, h, fill=None, line=None, lw=1, radius=None):
    shp = SH.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        E(x), E(y), E(w), E(h))
    if radius:
        shp.adjustments[0] = radius / min(w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = P(lw)
    shp.text_frame.text = ""
    no_shadow(shp)
    return shp


def textbox(x, y, w, h, parts, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
            wrap=True, tag=None):
    """parts: [(runs, size, weight, color, line_height, spc_em, space_before)]
       runs: str 또는 [(text, color|None, weight|None)]"""
    tb = SH.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    used = 0.0
    for i, (runs, size, weight, color, lh, spc, sb) in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = P(size * lh)
        if sb:
            p.space_before = P(sb)
        p.space_after = Pt(0)
        if isinstance(runs, str):
            runs = [(runs, None, None)]
        flat = "".join(t for t, _, _ in runs)
        if tag and wrap:
            n = est_lines(flat, size, w, weight, spc)
        else:
            n = 1
            if tag and text_w(flat, size, weight, spc) > w + 0.5:
                WARNINGS.append(f"[wide] {tag}: {text_w(flat, size, weight, spc):.0f}px "
                                f"> {w}px | {flat[:44]}")
        used += n * size * lh + (sb or 0)
        for txt, rcolor, rweight in runs:
            r = p.add_run()
            r.text = txt
            wt = rweight if rweight else weight
            name, bold = face(wt)
            r.font.name = name
            r.font.bold = bold
            r.font.size = P(size)
            r.font.color.rgb = rcolor if rcolor else color
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(int(round(size * 0.75 * spc * 100))))
            rPr.set('kern', '0')
            latin = rPr.find(qn('a:latin'))
            for t in ('a:ea', 'a:cs'):
                el = rPr.makeelement(qn(t), {'typeface': name})
                rPr.insert(list(rPr).index(latin) + 1, el)
    if tag and used > h + 0.5:
        WARNINGS.append(f"[overflow] {tag}: 내용 {used:.0f}px > 상자 {h}px")
    return tb


def line(x1, y1, x2, y2, color, w_px, cap_round=True, arrow=False, dash=None):
    c = SH.add_connector(MSO_CONNECTOR.STRAIGHT, E(x1), E(y1), E(x2), E(y2))
    c.line.color.rgb = color
    c.line.width = P(w_px)
    ln = c.line._get_or_add_ln()
    if cap_round:
        ln.set('cap', 'rnd')
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)
    if arrow:
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    no_shadow(c)
    return c


def oval(cx, cy, r, fill, line_c=None, lw=1):
    shp = SH.add_shape(MSO_SHAPE.OVAL, E(cx - r), E(cy - r), E(r * 2), E(r * 2))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_c is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_c
        shp.line.width = P(lw)
    shp.text_frame.text = ""
    no_shadow(shp)
    return shp


def picture(path, x, y, w, h):
    p = SH.add_picture(path, E(x), E(y), E(w), E(h))
    return p


# ── §6 브랜드 마크 ────────────────────────────────────────────────────────────
def mark(x, y, size, dark=False):
    """viewBox 100x100 기준 마크를 (x, y)에 size px 로 그린다."""
    s = size / 100.0
    mech = ON_DARK if dark else INK
    acc = ACCENT_D if dark else ACCENT

    def mx(v): return x + v * s
    def my(v): return y + v * s
    def mr(v): return v * s

    rect(mx(27), my(39), mr(40), mr(33), fill=mech, line=mech, lw=max(2.5 * s, 0.25),
         radius=mr(14))                                    # 몸체
    line(mx(47), my(39), mx(47), my(30), mech, 3.5 * s)     # 비콘 기둥
    oval(mx(47), my(28), mr(3.4), acc)                      # 비콘
    line(mx(65), my(47), mx(81), my(37), mech, 5.5 * s)     # 상완
    oval(mx(81), my(37), mr(3.5), mech)                     # 팔꿈치
    line(mx(81), my(37), mx(88), my(51), mech, 5.5 * s)     # 전완
    oval(mx(88), my(51), mr(4.2), acc)                      # 손목
    line(mx(88), my(51), mx(95), my(45), acc, 4 * s)        # 집게 위
    line(mx(88), my(51), mx(94), my(58), acc, 4 * s)        # 집게 아래
    oval(mx(38), my(74), mr(6.5), mech)                     # 바퀴
    oval(mx(56), my(74), mr(6.5), mech)                     # 바퀴
    oval(mx(50), my(55), mr(9), acc)                        # 눈
    oval(mx(53), my(53), mr(2.2), BLACK)                    # 동공


# ── §3 헤더 · 푸터 ────────────────────────────────────────────────────────────
def ground(dark=False):
    rect(0, 0, 1280, 720, fill=TILE1 if dark else GROUND)


def header(rail_text, title, subtitle="", dark=False):
    rect(64, 50, 6, 6, fill=ACCENT_D if dark else ACCENT)
    textbox(76, 44, 1140, 18,
            [(rail_text, 13, 600, ON_DARK_M if dark else INK_DIM, 1.0, 0.12, 0)],
            anchor=MSO_ANCHOR.MIDDLE, wrap=False, tag="rail")
    textbox(64, 70, 1152, 46,
            [(title, 36, 600, ON_DARK if dark else INK, 1.15, -0.02, 0)],
            anchor=MSO_ANCHOR.MIDDLE, wrap=False, tag="title")
    if subtitle:
        textbox(64, 121, 1152, 26,
                [(subtitle, 17, 400, ON_DARK_M if dark else INK_DIM, 1.45, -0.01, 0)],
                anchor=MSO_ANCHOR.MIDDLE, wrap=False, tag="subtitle")
    rect(64, 158, 1152, 1, fill=RULE_DARK if dark else RULE)


def footer(page, total=20, dark=False):
    mark(64, 676, 18, dark=dark)
    c = ON_DARK_M if dark else INK_FNT
    textbox(90, 676, 500, 20,
            [("HazardBot 개발완료보고서", 11, 600, c, 1.0, 0.06, 0)],
            anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    textbox(916, 676, 300, 20,
            [(f"{page:02d} / {total}", 11, 600, c, 1.0, 0.06, 0)],
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT, wrap=False)


# ── §11 컨텐츠 블록 ───────────────────────────────────────────────────────────
def field(x, y, w, h, radius=6):
    """diagram-field / image-mat 의 바탕 — 흰 면 + 헤어라인."""
    return rect(x, y, w, h, fill=CANVAS, line=RULE, lw=1, radius=radius)


def card(x, y, w, h, eyebrow=None, title=None, body=None, fill=CANVAS,
         accent_bar=False, pad=24, tag=""):
    """§11 card — 흰 면 + 헤어라인. accent_bar 면 card-accent (슬라이드당 1개)."""
    rect(x, y, w, h, fill=WASH if accent_bar else fill, line=RULE, lw=1, radius=6)
    if accent_bar:
        rect(x, y + 6, 3, h - 12, fill=ACCENT)
    iw = w - pad * 2
    cy = y + pad
    if eyebrow:
        textbox(x + pad, cy, iw, 15, [(eyebrow, 12, 600, INK_DIM, 1.2, 0.08, 0)],
                wrap=False, tag=tag + ".eyebrow")
        cy += 15 + 12
    if title:
        n = check(tag + ".title", title, 19, iw, 999, 1.35, 600, -0.015)
        textbox(x + pad, cy, iw, n * 26, [(title, 19, 600, INK, 1.35, -0.015, 0)],
                tag=tag + ".title")
        cy += n * 26 + 16
    if body:
        textbox(x + pad, cy, iw, y + h - pad - cy,
                [(body, 14, 400, INK_DIM, 1.6, 0, 0)], tag=tag + ".body")
    return cy


def rows(x, y, w, h, items, name_size=14, desc_size=12, divider=True,
         name_color=INK, desc_color=INK_DIM, max_slack=None, tag=""):
    """헤어라인으로 나뉜 목록 행 — 카드 안의 목록. items = [(name, desc)]

    행 높이는 내용량에 비례해 나누고 남는 여백만 균등 분배한다.

    max_slack 을 주면 행당 추가 여백을 그만큼으로 제한하고 나머지는 아래에 남긴다.
    항목 수가 카드마다 다를 때(팀 구성 등) 이걸 쓰지 않으면, 항목이 적은 카드의
    행들이 카드 높이만큼 벌어져 옆 카드와 리듬이 어긋난다."""
    n = len(items)
    need = []
    for nm, ds in items:
        hh = est_lines(nm, name_size, w, 600) * name_size * 1.5
        if ds:
            hh += 3 + est_lines(ds, desc_size, w) * desc_size * 1.5
        need.append(hh)
    slack = (h - (n - 1) - sum(need)) / n
    if max_slack is not None:
        slack = min(slack, max_slack)
    if slack < 0:
        WARNINGS.append(f"[overflow] {tag}: 행 내용 {sum(need):.0f}px > 영역 {h}px")
        slack = 0
    top = y
    for i, (nm, ds) in enumerate(items):
        ih = need[i] + slack
        if i and divider:
            rect(x, top - 1, w, 1, fill=RULE_SFT)
        parts = [(nm, name_size, 600, name_color, 1.5, 0, 0)]
        if ds:
            parts.append((ds, desc_size, 400, desc_color, 1.5, 0, 3))
        textbox(x, top, w, ih, parts, anchor=MSO_ANCHOR.MIDDLE,
                tag=f"{tag}.row{i}")
        top += ih + 1


def metric_strip(x, y, w, metrics, height=144):
    """§11 metric-strip — full(144)은 label + figure + qualifier, compact(96)은 label + figure."""
    rect(x, y, w, height, fill=PEARL, line=RULE, lw=1, radius=6)
    n = len(metrics)
    inner = w - 48
    cw = (inner - (n - 1)) / n
    for i, m in enumerate(metrics):
        label, fig = m[0], m[1]
        sub = m[2] if len(m) > 2 else None
        cx = x + 24 + i * (cw + 1)
        if i:
            rect(cx - 1, y + 24, 1, height - 48, fill=RULE_SFT)
        tx = cx + (24 if i else 0)
        tw = cw - (24 if i else 0) - (24 if i < n - 1 else 0)
        if height >= 144:
            textbox(tx, y + 28.8, tw, 15, [(label, 12, 600, INK_DIM, 1.2, 0.08, 0)],
                    wrap=False, tag=f"metric{i}.label")
            textbox(tx, y + 49.2, tw, 44, [(fig, 44, 700, ACCENT, 1.0, -0.02, 0)],
                    wrap=False, tag=f"metric{i}.fig")
            if sub:
                textbox(tx, y + 97.2, tw, 36, [(sub, 12, 400, INK_DIM, 1.5, 0, 0)],
                        tag=f"metric{i}.sub")
        else:
            textbox(tx, y + 20, tw, 15, [(label, 12, 600, INK_DIM, 1.2, 0.08, 0)],
                    wrap=False, tag=f"metric{i}.label")
            textbox(tx, y + 40, tw, 36, [(fig, 32, 700, ACCENT, 1.0, -0.02, 0)],
                    wrap=False, tag=f"metric{i}.fig")


def evidence_rail(x, y, w, heading, items, height=144):
    """§11 evidence-rail — 근거·출처를 본문 하단에 고정한다. items = [(label, text)]"""
    rect(x, y, w, height, fill=PEARL, line=RULE, lw=1, radius=6)
    textbox(x + 24, y + 22, w - 48, 15, [(heading, 12, 600, INK_DIM, 1.2, 0.08, 0)],
            wrap=False, tag="rail.heading")
    n = len(items)
    inner = w - 48
    cw = (inner - (n - 1)) / n
    top = y + 52
    ch = height - 52 - 22
    for i, (label, text) in enumerate(items):
        cx = x + 24 + i * (cw + 1)
        if i:
            rect(cx - 1, top, 1, ch, fill=RULE_SFT)
        tx = cx + (20 if i else 0)
        tw = cw - (20 if i else 0) - (20 if i < n - 1 else 0)
        textbox(tx, top, tw, ch,
                [(label, 12, 600, INK, 1.4, 0, 0),
                 (text, 11, 400, INK_DIM, 1.45, 0, 4)], tag=f"rail{i}")


def node(x, y, w, h, title, desc=None, active=False, ground_fill=GROUND,
         swatch=None, size=14, align=PP_ALIGN.LEFT, pad=12, tag=""):
    """§11 diagram-field 안의 노드. active 는 슬라이드당 1개."""
    rect(x, y, w, h, fill=WASH if active else ground_fill,
         line=ACCENT_F if active else INK_FNT, lw=1.5, radius=6)
    pad_l = 26 if swatch else pad
    if swatch:
        oval(x + 15, y + h / 2 - (7 if desc else 0), 5, swatch)
    parts = [(title, size, 600, INK, 1.45, 0, 0)]
    if desc:
        parts.append((desc, 12, 400, INK_DIM, 1.4, 0, 4))
    textbox(x + pad_l, y + 8, w - pad_l - pad, h - 16, parts, align=align,
            anchor=MSO_ANCHOR.MIDDLE, tag=tag or f"node.{title[:8]}")


def badge(x, y, text, kind, on_ground=False, h=20):
    fg, bg_w, bg_g = STATUS[kind]
    w = text_w(text, 12, 600, 0.08) + 20
    rect(x, y, w, h, fill=bg_g if on_ground else bg_w, radius=3)
    textbox(x, y, w, h, [(text, 12, 600, fg, 1.0, 0.08, 0)],
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, wrap=False)
    return w


def chip(x, y, text, size=13, fill=GROUND, color=INK, h=22):
    """§11 code-chip — 토픽·노드·파일명. 별도 고정폭 글꼴을 쓰지 않는다."""
    w = text_w(text, size, 400) + 14
    rect(x, y, w, h, fill=fill, radius=3)
    textbox(x, y, w, h, [(text, size, 400, color, 1.0, 0, 0)],
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, wrap=False)
    return w


def callout(x, y, w, h, parts, pad=24):
    """§11 callout — primary-wash 면 + 3px primary 좌측 바. 슬라이드당 1개."""
    rect(x, y, w, h, fill=WASH, line=None, radius=6)
    rect(x, y + 6, 3, h - 12, fill=ACCENT)
    textbox(x + pad + 3, y + pad, w - pad * 2 - 3, h - pad * 2, parts,
            anchor=MSO_ANCHOR.MIDDLE, tag="callout")


def image_mat(x, y, w, h, path, img_w, img_h, caption_parts, pad=12):
    """§10 image-mat — 흰 매트 + 헤어라인 + 캡션."""
    rect(x, y, w, h, fill=CANVAS, line=RULE, lw=1, radius=6)
    picture(path, x + pad, y + pad, img_w, img_h)
    cy = y + pad + img_h + 8
    textbox(x + pad, cy, w - pad * 2, y + h - pad - cy, caption_parts, tag="caption")


def arrow(x1, y1, x2, y2, color=INK_FNT, w=1.5):
    line(x1, y1, x2, y2, color, w, arrow=True)


def elbow(pts, color=INK_FNT, w=1.5, arrow_last=True):
    """꺾인 연결선 — 마지막 구간에만 화살촉."""
    for i in range(len(pts) - 1):
        (a, b), (c, d) = pts[i], pts[i + 1]
        line(a, b, c, d, color, w, arrow=(arrow_last and i == len(pts) - 2))


def label(x, y, w, text, size=12, color=INK_DIM, align=PP_ALIGN.LEFT,
          weight=400, lh=1.4, spc=0.0, h=None, anchor=MSO_ANCHOR.TOP, tag=None):
    return textbox(x, y, w, h if h else size * lh + 2,
                   [(text, size, weight, color, lh, spc, 0)],
                   align=align, anchor=anchor, wrap=(h is not None), tag=tag)
